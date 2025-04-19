import os
import shutil
import random
from pptx import Presentation
import google.generativeai as genai
import webbrowser
import pyautogui
import time

# Mock Speech module if unavailable
try:
    from Speech import TextToSpeech as speak
except ImportError:
    def speak(text):
        print(f"Speech: {text}")

# Configure Gemini API with secure key input
try:
    api_key = os.getenv("GEMINI_API_KEY")
    genai.configure(api_key=api_key)
except Exception as e:
    speak(f"Failed to configure Gemini API: {str(e)}")
    exit(1)

# Use available model
gen_model = genai.GenerativeModel(model_name="gemini-1.5-flash")

PRESENTATION_DIR = "presentations"
DESIGNS_DIR = "designs"

topic = "presentation on migration of birds"
os.makedirs(PRESENTATION_DIR, exist_ok=True)
os.makedirs(DESIGNS_DIR, exist_ok=True)

def clean_old_ppts():
    try:
        for file in os.listdir(PRESENTATION_DIR):
            if file.endswith(".pptx"):
                os.remove(os.path.join(PRESENTATION_DIR, file))
        speak("I have deleted the older files in the directory for your ease.")
    except Exception as e:
        speak(f"Error cleaning old presentations: {str(e)}")

filename = os.path.join(PRESENTATION_DIR, f"{topic}.pptx")

def get_random_template():
    try:
        templates = [os.path.join(DESIGNS_DIR, f) for f in os.listdir(DESIGNS_DIR) if f.endswith(".pptx")]
        if not templates:
            speak("No templates found in designs directory. Using default presentation.")
            return None
        return random.choice(templates)
    except Exception as e:
        speak(f"Error selecting template: {str(e)}")
        return None

def generate_slide_content(topic, slide_num):
    prompt = (
    "This AI is a part of the greater artificial intelligence called Jarvis. "
    "It must generate structured PowerPoint slides for a presentation on '{topic}', with unique titles and 3-5 bullet points per slide. "
    "The AI must maintain logical flow across slides and ensure that each slide contributes to the overall structure. "
    "It should include the following: an introduction slide, an agenda slide, 6-8 content slides (each covering a distinct subtopic), and a conclusion slide. "
    "Each slide should have a clear and meaningful title related to its content. "
    "The AI must not generate any blank or repetitive slides. "
    "Do not use any formatting symbols (no asterisks, no bold markers, no underscores). "
    "Only provide raw text. Generate slide number {slide_num} in sequence for the topic '{topic}'."
).format(slide_num=slide_num, topic=topic)

     
    try:
        response = gen_model.generate_content(prompt)
        if not response.text:
            speak("Failed to generate content for slide.")
            return None
        
        # Safety check for content
        if hasattr(response, 'prompt_feedback') and response.prompt_feedback:
            for rating in response.prompt_feedback.safety_ratings:
                if rating.category == 'HARM_CATEGORY_SEXUALLY_EXPLICIT' and rating.probability in ['MEDIUM', 'HIGH']:
                    speak("Content flagged for harmful content. Skipping slide.")
                    return None

        clean_text = response.text.replace("**", "").replace("*", "").replace("-", "").strip()
        return [line.strip() for line in clean_text.split("\n") if line.strip()]
    except Exception as e:
        print(f"Error generating slide content: {str(e)}")
        return None

def create_ppt(topic, num_slides=10):
    try:
        clean_old_ppts()
        template_path = get_random_template()
        speak("I will let you know when the ppt is created")
        prs = Presentation(template_path) if template_path else Presentation()

        for i in range(1, num_slides + 1):
            content = generate_slide_content(topic, i)
            slide_layout = prs.slide_layouts[1]  # Title + Content layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content_box = slide.shapes.placeholders[1]

            if content:
                title.text = content[0][:50]  # Truncate title if too long
                for point in content[1:]:
                    p = content_box.text_frame.add_paragraph()
                    p.text = point

        prs.save(filename)
        speak(f"Sir, the ppt based on your topic is created.")
    except Exception as e:
        print(f"Error creating presentation: {str(e)}")

def makepresentation():
    try:
        if not os.path.exists(filename):
            speak("Presentation file not found.")
            return
        webbrowser.open(filename)
        speak("Opening the file, sir.")
        time.sleep(5)  # Increased delay for application to open
    except Exception as e:
        print(f"Error opening presentation: {str(e)}")

def start_slideshow():
    try:
        pyautogui.press("f5")
        speak("Starting the slideshow, sir.")
        time.sleep(1)
    except Exception as e:
        print(f"Error starting slideshow: {str(e)}")

def next_slide():
    try:
        pyautogui.press("right")
        speak("Next slide.")
        time.sleep(1)
    except Exception as e:
        print(f"Error navigating to next slide: {str(e)}")

def previous_slide():
    try:
        pyautogui.press("left")
        speak("Previous slide.")
        time.sleep(1)
    except Exception as e:
        print(f"Error navigating to previous slide: {str(e)}")

# Execute the presentation workflow
create_ppt(topic)
makepresentation()
start_slideshow()
next_slide()
next_slide()
next_slide()
previous_slide()
previous_slide()
previous_slide()